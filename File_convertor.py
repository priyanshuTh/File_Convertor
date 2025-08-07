import os
import pandas as pd
import json
import subprocess
import tempfile
import shutil
from fpdf import FPDF
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation
from PyPDF2 import PdfReader

class FileConverter:
    @staticmethod
    def _validate_extension(file_path, valid_exts, file_desc="file"):
        """Ensure the file has one of the valid extensions."""
        ext = os.path.splitext(file_path)[1].lower()
        if ext not in valid_exts:
            raise ValueError(f"Invalid extension for {file_desc}: '{ext}'. Expected one of {valid_exts}.")

    @staticmethod
    def csv_to_excel(csv_file, excel_file):
        FileConverter._validate_extension(csv_file, ['.csv'], 'CSV input')
        FileConverter._validate_extension(excel_file, ['.xls', '.xlsx'], 'Excel output')
        try:
            df = pd.read_csv(csv_file)
            df.to_excel(excel_file, index=False, engine='openpyxl')
            print(f"Converted {csv_file} to {excel_file}.")
        except Exception as e:
            print(f"Error converting CSV to Excel: {e}")

    @staticmethod
    def excel_to_csv(excel_file, csv_file):
        FileConverter._validate_extension(excel_file, ['.xls', '.xlsx'], 'Excel input')
        FileConverter._validate_extension(csv_file, ['.csv'], 'CSV output')
        try:
            df = pd.read_excel(excel_file, engine='openpyxl')
            df.to_csv(csv_file, index=False)
            print(f"Converted {excel_file} to {csv_file}.")
        except Exception as e:
            print(f"Error converting Excel to CSV: {e}")

    @staticmethod
    def text_to_json(text_file, json_file):
        FileConverter._validate_extension(text_file, ['.txt', '.text'], 'Text input')
        FileConverter._validate_extension(json_file, ['.json'], 'JSON output')
        try:
            with open(text_file, 'r', encoding='utf-8') as f:
                lines = [line.rstrip('\n') for line in f]
            data = {"lines": lines}
            with open(json_file, 'w', encoding='utf-8') as f:
                json.dump(data, f, indent=4)
            print(f"Converted {text_file} to {json_file}.")
        except Exception as e:
            print(f"Error converting text to JSON: {e}")

    @staticmethod
    def json_to_text(json_file, text_file):
        FileConverter._validate_extension(json_file, ['.json'], 'JSON input')
        FileConverter._validate_extension(text_file, ['.txt'], 'Text output')
        try:
            with open(json_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
            lines = data.get("lines", [])
            with open(text_file, 'w', encoding='utf-8') as f:
                f.write("\n".join(lines))
            print(f"Converted {json_file} to {text_file}.")
        except Exception as e:
            print(f"Error converting JSON to text: {e}")

    @staticmethod
    def text_to_csv(text_file, csv_file):
        FileConverter._validate_extension(text_file, ['.txt'], 'Text input')
        FileConverter._validate_extension(csv_file, ['.csv'], 'CSV output')
        try:
            with open(text_file, 'r', encoding='utf-8') as f:
                lines = [line.rstrip('\n') for line in f]
            df = pd.DataFrame({"lines": lines})
            df.to_csv(csv_file, index=False)
            print(f"Converted {text_file} to {csv_file}.")
        except Exception as e:
            print(f"Error converting text to CSV: {e}")

    @staticmethod
    def html_to_pdf(html_file, pdf_file):
        FileConverter._validate_extension(html_file, ['.html', '.htm'], 'HTML input')
        FileConverter._validate_extension(pdf_file, ['.pdf'], 'PDF output')
        try:
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            with open(html_file, 'r', encoding='utf-8') as f:
                for line in f:
                    pdf.multi_cell(0, 10, line.strip())
            pdf.output(pdf_file)
            print(f"Converted {html_file} to {pdf_file}.")
        except Exception as e:
            print(f"Error converting HTML to PDF: {e}")

    @staticmethod
    def pdf_to_html(pdf_file, html_file):
        FileConverter._validate_extension(pdf_file, ['.pdf'], 'PDF input')
        FileConverter._validate_extension(html_file, ['.html'], 'HTML output')
        try:
            reader = PdfReader(pdf_file)
            with open(html_file, 'w', encoding='utf-8') as f:
                f.write("<html><body>")
                for page in reader.pages:
                    text = page.extract_text() or ''
                    f.write(f"<p>{text.replace('\n', '<br>')}</p>")
                f.write("</body></html>")
            print(f"Converted {pdf_file} to {html_file}.")
        except Exception as e:
            print(f"Error converting PDF to HTML: {e}")

    @staticmethod
    def pdf_to_ppt(pdf_file, ppt_file):
        FileConverter._validate_extension(pdf_file, ['.pdf'], 'PDF input')
        FileConverter._validate_extension(ppt_file, ['.pptx'], 'PPTX output')
        try:
            presentation = Presentation()
            for img in convert_from_path(pdf_file):
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                img.save(tmp.name, 'PNG')
                slide.shapes.add_picture(tmp.name, 0, 0, width=presentation.slide_width)
                tmp.close()
                os.unlink(tmp.name)
            presentation.save(ppt_file)
            print(f"Converted {pdf_file} to {ppt_file}.")
        except Exception as e:
            print(f"Error converting PDF to PPT: {e}")

    @staticmethod
    def ppt_to_pdf(ppt_file, pdf_file):
        FileConverter._validate_extension(ppt_file, ['.ppt', '.pptx'], 'PPT input')
        FileConverter._validate_extension(pdf_file, ['.pdf'], 'PDF output')
        try:
            subprocess.run([
                'libreoffice', '--headless', '--convert-to', 'pdf', '--outdir',
                os.path.dirname(pdf_file) or '.', ppt_file
            ], check=True)
            orig = os.path.join(os.path.dirname(pdf_file),
                                os.path.splitext(os.path.basename(ppt_file))[0] + '.pdf')
            if orig != pdf_file:
                shutil.move(orig, pdf_file)
            print(f"Converted {ppt_file} to {pdf_file}.")
        except Exception as e:
            print(f"Error converting PPT to PDF: {e}")

    @staticmethod
    def pdf_to_jpg(pdf_file, output_folder):
        FileConverter._validate_extension(pdf_file, ['.pdf'], 'PDF input')
        try:
            os.makedirs(output_folder, exist_ok=True)
            for i, img in enumerate(convert_from_path(pdf_file)):
                img.save(os.path.join(output_folder, f'page_{i+1}.jpg'), 'JPEG')
            print(f"Converted {pdf_file} to JPG images in {output_folder}.")
        except Exception as e:
            print(f"Error converting PDF to JPG: {e}")

    @staticmethod
    def jpg_to_pdf(jpg_file, pdf_file):
        FileConverter._validate_extension(jpg_file, ['.jpg', '.jpeg'], 'JPG input')
        FileConverter._validate_extension(pdf_file, ['.pdf'], 'PDF output')
        try:
            img = Image.open(jpg_file)
            img.convert('RGB').save(pdf_file)
            print(f"Converted {jpg_file} to {pdf_file}.")
        except Exception as e:
            print(f"Error converting JPG to PDF: {e}")

    @staticmethod
    def png_to_pdf(png_file, pdf_file):
        FileConverter._validate_extension(png_file, ['.png'], 'PNG input')
        FileConverter._validate_extension(pdf_file, ['.pdf'], 'PDF output')
        try:
            img = Image.open(png_file)
            img.convert('RGB').save(pdf_file)
            print(f"Converted {png_file} to {pdf_file}.")
        except Exception as e:
            print(f"Error converting PNG to PDF: {e}")

    @staticmethod
    def pdf_to_png(pdf_file, output_folder):
        FileConverter._validate_extension(pdf_file, ['.pdf'], 'PDF input')
        try:
            os.makedirs(output_folder, exist_ok=True)
            for i, img in enumerate(convert_from_path(pdf_file)):
                img.save(os.path.join(output_folder, f'page_{i+1}.png'), 'PNG')
            print(f"Converted {pdf_file} to PNG images in {output_folder}.")
        except Exception as e:
            print(f"Error converting PDF to PNG: {e}")

    @staticmethod
    def gif_to_pdf(gif_file, pdf_file):
        FileConverter._validate_extension(gif_file, ['.gif'], 'GIF input')
        FileConverter._validate_extension(pdf_file, ['.pdf'], 'PDF output')
        try:
            img = Image.open(gif_file)
            frames = [img.convert('RGB') for _ in range(getattr(img, 'n_frames', 1))]
            frames[0].save(pdf_file, save_all=True, append_images=frames[1:])
            print(f"Converted {gif_file} to {pdf_file}.")
        except Exception as e:
            print(f"Error converting GIF to PDF: {e}")

    @staticmethod
    def pdf_to_gif(pdf_file, gif_file):
        FileConverter._validate_extension(pdf_file, ['.pdf'], 'PDF input')
        FileConverter._validate_extension(gif_file, ['.gif'], 'GIF output')
        try:
            frames = [img.convert('RGB') for img in convert_from_path(pdf_file)]
            frames[0].save(gif_file, save_all=True, append_images=frames[1:], duration=500, loop=0)
            print(f"Converted {pdf_file} to {gif_file}.")
        except Exception as e:
            print(f"Error converting PDF to GIF: {e}")

    @staticmethod
    def main():
        options = {
            '1': ('CSV to Excel', FileConverter.csv_to_excel),
            '2': ('Excel to CSV', FileConverter.excel_to_csv),
            '3': ('Text to JSON', FileConverter.text_to_json),
            '4': ('JSON to Text', FileConverter.json_to_text),
            '5': ('Text to CSV', FileConverter.text_to_csv),
            '6': ('HTML to PDF', FileConverter.html_to_pdf),
            '7': ('PDF to HTML', FileConverter.pdf_to_html),
            '8': ('PDF to PPT', FileConverter.pdf_to_ppt),
            '9': ('PPT to PDF', FileConverter.ppt_to_pdf),
            '10': ('PDF to JPG', FileConverter.pdf_to_jpg),
            '11': ('JPG to PDF', FileConverter.jpg_to_pdf),
            '12': ('PNG to PDF', FileConverter.png_to_pdf),
            '13': ('PDF to PNG', FileConverter.pdf_to_png),
            '14': ('GIF to PDF', FileConverter.gif_to_pdf),
            '15': ('PDF to GIF', FileConverter.pdf_to_gif)
        }
        while True:
            print("\nOptions:")
            for k, (desc, _) in options.items():
                print(f"{k}. {desc}")
            print("0. Exit")
            choice = input("Choice: ").strip().strip('"').strip("'")
            if choice == '0':
                break
            if choice in options:
                _, func = options[choice]
                src = input("Source path: ").strip().strip('"').strip("'")
                if choice in ['10', '13']:
                    out = input("Output folder: ").strip().strip('"').strip("'")
                    os.makedirs(out, exist_ok=True)
                    func(src, out)
                else:
                    dst = input("Destination path: ").strip().strip('"').strip("'")
                    func(src, dst)
            else:
                print("Invalid choice.")

if __name__ == '__main__':
    FileConverter.main()
